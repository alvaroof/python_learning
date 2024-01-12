# -*- coding: utf-8 -*-
from pptx import Presentation


class PPTProcessor:
    @staticmethod
    def process_ppt(file_path):
        """Processes the PowerPoint file and extracts text from each slide."""
        prs = Presentation(file_path)
        content = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content.append(shape.text)
        return " ".join(content)


# Test the processor (you can remove this test in your actual application)
if __name__ == "__main__":
    ppt_content = PPTProcessor.process_ppt("path_to_your_ppt.pptx")
    print(ppt_content)
